"""Generate the presentation figures and assemble a polished PowerPoint deck.

Produces:
  slides/figures/*.png                 (charts built from the saved artifacts)
  slides/Multi-LexSum_Presentation.pptx

Run:  python src/make_slides.py
"""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

import config as C

FIG = C.ROOT / "slides" / "figures"
FIG.mkdir(parents=True, exist_ok=True)
DECK = C.ROOT / "slides" / "Multi-LexSum_Presentation.pptx"

METRICS = json.load(open(C.MODELS_DIR / "metrics.json"))
TOPF = json.load(open(C.MODELS_DIR / "top_features.json"))

# Real ROUGE F1 values measured in notebooks/02_summarization.ipynb
ROUGE = {
    "rouge1": {"tiny (1)": 0.151, "short (4)": 0.265, "long (12)": 0.340},
    "rouge2": {"tiny (1)": 0.020, "short (4)": 0.060, "long (12)": 0.099},
    "rougeL": {"tiny (1)": 0.106, "short (4)": 0.148, "long (12)": 0.156},
}
ROUGE_METHODS = {"lexrank": 0.152, "textrank": 0.108, "lsa": 0.138}  # rougeL, short

# ---------------------------------------------------------------------------
# Shared brand palette (hex strings; the deck reuses the same tones as RGB)
# ---------------------------------------------------------------------------
INK      = "#16263A"   # near-black navy, body text
PRIMARY  = "#2D5BA8"   # brand blue
TEAL     = "#2A9D8F"   # secondary
CORAL    = "#E07A5F"   # warm accent / highlight
GOLD     = "#E9B949"
MUTED    = "#6B7A8D"
LIGHTBG  = "#F4F7FB"   # slide background
CARD     = "#FFFFFF"
SERIES   = [PRIMARY, CORAL, TEAL, GOLD, "#8E6FB6", "#5AA9E6"]


def set_style():
    plt.rcParams.update({
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "savefig.facecolor": "white",
        "axes.edgecolor": "#C9D3DF",
        "axes.linewidth": 0.9,
        "axes.grid": True,
        "grid.color": "#E5EBF2",
        "grid.linewidth": 0.8,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.titlesize": 14,
        "axes.titleweight": "bold",
        "axes.titlecolor": INK,
        "axes.labelcolor": MUTED,
        "axes.labelsize": 11,
        "xtick.color": MUTED,
        "ytick.color": MUTED,
        "font.size": 11,
        "legend.frameon": False,
    })


# ---------------------------------------------------------------------------
# Figures
# ---------------------------------------------------------------------------
def _save(fig, name):
    fig.savefig(FIG / name, dpi=160, bbox_inches="tight")
    plt.close(fig)


def fig_targets():
    df = pd.read_parquet(C.SUBSET_PARQUET)
    fig, ax = plt.subplots(1, 2, figsize=(11.5, 4.2),
                           gridspec_kw={"wspace": 0.4})
    vc = df.class_action_sought.value_counts().reindex(["No", "Yes"])
    bars = ax[0].bar(vc.index, vc.values, color=[PRIMARY, CORAL], width=0.6)
    ax[0].bar_label(bars, padding=3, color=INK, fontsize=11)
    ax[0].set_title("class_action_sought  (binary target)")
    ax[0].set_ylabel("# cases"); ax[0].grid(axis="x", visible=False)
    ct = df.case_type.value_counts()
    ax[1].barh(ct.index, ct.values, color=TEAL)
    ax[1].invert_yaxis(); ax[1].grid(axis="y", visible=False)
    ax[1].set_title("case_type  (12-class target)")
    _save(fig, "targets.png")


def fig_rouge():
    fig, ax = plt.subplots(1, 2, figsize=(11, 4.2))
    rg = pd.DataFrame(ROUGE)
    rg.plot.bar(ax=ax[0], color=[PRIMARY, CORAL, TEAL], rot=0, width=0.75,
                edgecolor="white")
    ax[0].set_title("ROUGE F1 vs expert summaries, by tier")
    ax[0].set_ylabel("F1"); ax[0].legend(title="", ncol=3, loc="upper left")
    ax[0].grid(axis="x", visible=False)
    s = pd.Series(ROUGE_METHODS)
    b = ax[1].bar(s.index, s.values, color=[PRIMARY, MUTED, TEAL], width=0.6)
    ax[1].bar_label(b, fmt="%.3f", padding=3, color=INK)
    ax[1].set_title("ROUGE-L by summariser (short tier)")
    ax[1].set_ylabel("rougeL F1"); ax[1].grid(axis="x", visible=False)
    _save(fig, "rouge.png")


def fig_leaderboard():
    fig, ax = plt.subplots(1, 2, figsize=(12, 4.2),
                           gridspec_kw={"wspace": 0.55})
    for i, target in enumerate(["class_action_sought", "case_type"]):
        d = METRICS[target]["models"]
        s = pd.Series({m: d[m]["accuracy"] for m in d}).sort_values()
        colors = [CORAL if m == s.index[-1] else PRIMARY for m in s.index]
        ax[i].barh(s.index, s.values, color=colors)
        ax[i].set_xlim(0, 1.18); ax[i].set_title(f"Test accuracy — {target}")
        ax[i].grid(axis="y", visible=False)
        for y, v in enumerate(s.values):
            ax[i].text(v + 0.02, y, f"{v:.3f}", va="center", fontsize=10,
                       color=INK, fontweight="bold")
    _save(fig, "leaderboard.png")


def fig_perclass():
    d = METRICS["case_type"]["models"]["linear_svm"]["per_class_accuracy"]
    s = pd.Series(d).sort_values()
    fig, ax = plt.subplots(figsize=(8, 5))
    colors = [CORAL if v < 0.8 else TEAL for v in s.values]
    ax.barh(s.index, s.values, color=colors)
    ax.set_xlim(0, 1.05); ax.grid(axis="y", visible=False)
    ax.set_title("Per-class accuracy — case_type (Linear SVM)")
    for y, v in enumerate(s.values):
        ax.text(min(v + 0.015, 0.92), y, f"{v:.2f}", va="center", fontsize=9, color=INK)
    _save(fig, "perclass.png")


def fig_confusion():
    cm = np.array(METRICS["case_type"]["models"]["linear_svm"]["confusion"], float)
    labs = METRICS["case_type"]["labels"]
    cmn = cm / cm.sum(axis=1, keepdims=True)
    fig, ax = plt.subplots(figsize=(8.5, 7.5))
    im = ax.imshow(cmn, cmap="Blues", vmin=0, vmax=1)
    cb = plt.colorbar(im, fraction=0.046, pad=0.04); cb.outline.set_visible(False)
    ax.set_xticks(range(len(labs))); ax.set_xticklabels(labs, rotation=90, fontsize=7)
    ax.set_yticks(range(len(labs))); ax.set_yticklabels(labs, fontsize=7)
    ax.set_xlabel("predicted"); ax.set_ylabel("true")
    ax.set_title("case_type confusion matrix (row-normalised)")
    ax.grid(False)
    for i in range(len(labs)):
        for j in range(len(labs)):
            v = cmn[i, j]
            if v >= 0.08:
                ax.text(j, i, f"{v:.0%}", ha="center", va="center", fontsize=6,
                        color="white" if v > 0.5 else INK)
    _save(fig, "confusion.png")


def fig_wordclouds():
    from wordcloud import WordCloud

    def cloud(ax, words, title, cmap):
        freqs = {w: len(words) - j for j, w in enumerate(words)}
        wc = WordCloud(width=520, height=300, background_color="white",
                       colormap=cmap, prefer_horizontal=0.95
                       ).generate_from_frequencies(freqs)
        ax.imshow(wc); ax.axis("off")
        ax.set_title(title, fontsize=12, color=INK, fontweight="bold")

    fig, ax = plt.subplots(1, 2, figsize=(11, 3.6))
    cloud(ax[0], TOPF["class_action_sought"]["logreg"]["Yes"],
          "class_action_sought = Yes", "YlOrRd")
    cloud(ax[1], TOPF["class_action_sought"]["logreg"]["No"],
          "class_action_sought = No", "Blues")
    _save(fig, "wc_class_action.png")

    show = ["Equal Employment", "Prison Conditions",
            "Immigration and/or the Border", "Policing"]
    cmaps = ["Greens", "Purples", "Oranges", "GnBu"]
    fig, ax = plt.subplots(1, 4, figsize=(14, 3.1)); ax = ax.ravel()
    for i, cls in enumerate(show):
        cloud(ax[i], TOPF["case_type"]["logreg"].get(cls, []), cls, cmaps[i])
    _save(fig, "wc_case_type.png")


def fig_lstm():
    fig, ax = plt.subplots(1, 2, figsize=(11, 4))
    for i, target in enumerate(["class_action_sought", "case_type"]):
        h = METRICS[target]["models"]["lstm"]["history"]
        ep = range(1, len(h["accuracy"]) + 1)
        ax[i].plot(ep, h["accuracy"], "o-", color=PRIMARY, lw=2.2, label="train")
        ax[i].plot(ep, h["val_accuracy"], "s--", color=CORAL, lw=2.2, label="val")
        ax[i].set_title(f"Bi-LSTM accuracy / epoch — {target}")
        ax[i].set_xlabel("epoch"); ax[i].set_ylabel("accuracy")
        ax[i].set_ylim(0, 1); ax[i].legend(loc="lower right")
    _save(fig, "lstm.png")


def generate_figures():
    set_style()
    print("Generating figures ...")
    fig_targets(); fig_rouge(); fig_leaderboard(); fig_perclass()
    fig_confusion(); fig_wordclouds(); fig_lstm()
    print(f"  saved figures to {FIG}")


# ---------------------------------------------------------------------------
# Deck
# ---------------------------------------------------------------------------
def build_deck():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    def C_(hexv):  # hex -> RGBColor
        return RGBColor.from_string(hexv.lstrip("#"))

    cINK, cPRIM, cTEAL, cCORAL, cGOLD = (C_(INK), C_(PRIMARY), C_(TEAL),
                                         C_(CORAL), C_(GOLD))
    cMUTED, cBG, cWHITE = C_(MUTED), C_(LIGHTBG), RGBColor(255, 255, 255)
    cRULE = C_("#D8E0EA")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]
    state = {"n": 0}

    def rect(slide, shape, l, t, w, h, fill=None, line=None, line_w=None):
        sp = slide.shapes.add_shape(shape, l, t, w, h)
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            if line_w: sp.line.width = line_w
        sp.shadow.inherit = False
        return sp

    def tbox(slide, l, t, w, h, anchor=None):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        if anchor is not None: tf.vertical_anchor = anchor
        return tf

    def run(p, text, size, color, bold=False, italic=False, font="Calibri"):
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
        return r

    def base(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = cBG

    def footer(slide):
        state["n"] += 1
        tf = tbox(slide, Inches(0.55), SH - Inches(0.5), Inches(8), Inches(0.35))
        run(tf.paragraphs[0], "Multi-LexSum  ·  Text Analytics Option 2", 10, cMUTED)
        tf2 = tbox(slide, SW - Inches(1.4), SH - Inches(0.5), Inches(0.9), Inches(0.35))
        p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        run(p, f"{state['n']:02d}", 10, cMUTED, bold=True)

    def content_slide(title, subtitle=None):
        s = prs.slides.add_slide(BLANK); base(s)
        # left accent ribbon
        rect(s, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH, fill=cPRIM)
        # title
        tf = tbox(s, Inches(0.6), Inches(0.42), SW - Inches(1.1), Inches(1.0))
        run(tf.paragraphs[0], title, 28, cINK, bold=True)
        if subtitle:
            p = tf.add_paragraph(); run(p, subtitle, 14.5, cPRIM)
        # thin rule under header
        rect(s, MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.5), Inches(3.0),
             Pt(3), fill=cCORAL)
        footer(s)
        return s

    def bullets(s, items, left=Inches(0.7), top=Inches(1.75),
                width=None, height=None, size=18):
        width = width or SW - Inches(1.4)
        height = height or SH - Inches(2.4)
        tf = tbox(s, left, top, width, height)
        first = True
        for it in items:
            level = 0
            if isinstance(it, tuple):
                it, level = it
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(7); p.line_spacing = 1.08; p.level = level
            marker = "▪  " if level == 0 else "–  "
            run(p, marker, size - 2 * level, cCORAL if level == 0 else cTEAL, bold=True)
            run(p, it, size - 2 * level, cINK if level == 0 else cMUTED)
        return tf

    def img(s, path, l, t, w=None, h=None):
        if Path(path).exists():
            return s.shapes.add_picture(str(path), l, t, width=w, height=h)

    def stat_card(s, l, t, w, h, value, label, color):
        card = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=cWHITE,
                    line=cRULE, line_w=Pt(1))
        bar = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.12), fill=color)
        tf = card.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, value, 34, color, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        run(p2, label, 12.5, cMUTED)

    def notes(s, text):
        s.notes_slide.notes_text_frame.text = text

    # ================= Slide 1 — Title =================
    s = prs.slides.add_slide(BLANK); base(s)
    rect(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(3.05), fill=cINK)
    rect(s, MSO_SHAPE.RECTANGLE, 0, Inches(3.05), SW, Inches(0.16), fill=cCORAL)
    # accent blocks
    rect(s, MSO_SHAPE.RECTANGLE, SW - Inches(2.2), 0, Inches(0.6), Inches(3.05), fill=cPRIM)
    rect(s, MSO_SHAPE.RECTANGLE, SW - Inches(1.5), 0, Inches(0.6), Inches(3.05), fill=cTEAL)
    tf = tbox(s, Inches(0.85), Inches(0.7), SW - Inches(3.6), Inches(2.2))
    run(tf.paragraphs[0], "Summarising & Classifying", 38, cWHITE, bold=True)
    p = tf.add_paragraph(); run(p, "Civil-Rights Lawsuits", 38, cWHITE, bold=True)
    p = tf.add_paragraph(); run(p, "An NLP pipeline over the Multi-LexSum dataset",
                                19, C_("#BDD0E8"))
    tf2 = tbox(s, Inches(0.85), Inches(3.55), SW - Inches(1.7), Inches(3.2))
    for txt, sz, bold, col in [
        ("Text Analytics — Final Project (Option 2)", 21, True, cINK),
        ("Northwestern University · Spring 2025", 16, False, cMUTED),
        ("Team:  [ your names here ]", 16, False, cMUTED),
    ]:
        p = tf2.add_paragraph(); p.space_after = Pt(4)
        run(p, txt, sz, col, bold=bold)
    # pill summary
    pill = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.7),
                Inches(11.0), Inches(0.95), fill=cWHITE, line=cRULE, line_w=Pt(1))
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = pill.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    run(pp, "3-tier summaries  ·  class-action & case-type prediction  ·  live tool",
        16, cPRIM, bold=True)
    notes(s, "Hook: Multi-LexSum has 9,280 expert-written summaries of civil-rights "
             "cases whose source documents can exceed 200 pages. Our question: can we "
             "automatically reproduce those summaries AND predict case attributes? "
             "Today: data engineering, an extractive summariser, four classifiers on "
             "two targets, and a live Streamlit tool. (~10 minutes.)")

    # ================= Slide 2 — Dataset & task =================
    s = content_slide("The dataset & the task",
                      "Multi-LexSum (allenai) — expert civil-rights case summaries")
    bullets(s, [
        "9,280 civil-rights cases; source documents often exceed 200 pages.",
        "Each case ships THREE expert summaries at different granularity:",
        ("long (multi-paragraph)  ·  short (one paragraph)  ·  tiny (one sentence)", 1),
        "Rich metadata — two fields are our prediction targets:",
        ("class_action_sought  →  binary (Yes / No)", 1),
        ("case_type  →  multi-class (Prison Conditions, Equal Employment, …)", 1),
        "Deliverables: cleaning → summariser → 4 models × 2 targets → interactive tool.",
    ])
    notes(s, "Multi-LexSum is unusual: expert lawyers/law-students wrote the summaries "
             "to strict guidelines, reviewed by a second expert — a clean, high-quality "
             "multi-document summarisation benchmark. We use it for BOTH summarisation "
             "and two classification tasks. Stress the three-tier granularity — that's "
             "what our summariser must reproduce.")

    # ================= Slide 3 — Data engineering + EDA =================
    s = content_slide("Data engineering & EDA",
                      "Recovering full text from a 2.2 GB source file; a balanced subset")
    bullets(s, [
        "Full text lives in one ~2.2 GB JSON object with invalid NaN literals "
        "→ loaded with Python json, pulled per case.",
        "Balanced subset: 1,915 cases · 12 case types (≤180/type) · official "
        "train / dev / test preserved.",
    ], top=Inches(1.7), width=Inches(5.0), height=Inches(2.6), size=15)
    stat_card(s, Inches(0.7), Inches(4.55), Inches(2.4), Inches(1.3), "1,915", "cases", cPRIM)
    stat_card(s, Inches(3.25), Inches(4.55), Inches(2.4), Inches(1.3), "12", "case types", cTEAL)
    img(s, FIG / "targets.png", Inches(6.05), Inches(2.05), w=Inches(6.9))
    notes(s, "The engineering challenge: the 2.2 GB sources file isn't valid JSON (bare "
             "NaN), so streaming parsers like ijson choke — Python's json accepts NaN, "
             "so we load once and extract. We model a balanced subset for class balance "
             "and runtime; it scales to the full data via caps in config.py.")

    # ================= Slide 4 — Cleaning =================
    s = content_slide("Cleaning the legal text",
                      "Scanned, OCR'd court documents are noisy")
    bullets(s, [
        "Removed court page-headers, form-feeds, ECF/Doc stamps, lone page numbers.",
        "Repaired broken unicode (§, en/em dashes, the “��” mojibake).",
        "Second pass for modelling: lower-case + remove English AND legal stop-words.",
    ], top=Inches(1.7), height=Inches(1.7), size=16)
    # RAW card
    raw = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.75),
               Inches(11.9), Inches(1.35), fill=cWHITE, line=C_("#F0C9C0"), line_w=Pt(1.5))
    tf = raw.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2)
    run(tf.paragraphs[0], "RAW   ", 13, cCORAL, bold=True)
    run(tf.paragraphs[0], "Case 1:19-cv-01075-ERK-PK Document 134 Filed 05/25/21 Page 1 "
        "of 27 PageID #: 1374 … Plaintiffs brought this “FTCA” action under 42 U.S.C. "
        "�� 1983 – see ECF No. 29 at 1.", 12.5, cINK)
    # arrow
    ar = rect(s, MSO_SHAPE.DOWN_ARROW, Inches(6.35), Inches(5.18), Inches(0.6),
              Inches(0.34), fill=cPRIM)
    # CLEAN card
    cln = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.62),
               Inches(11.9), Inches(1.05), fill=cWHITE, line=C_("#BFE0C8"), line_w=Pt(1.5))
    tf = cln.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2)
    run(tf.paragraphs[0], "CLEAN   ", 13, cTEAL, bold=True)
    run(tf.paragraphs[0], 'Plaintiffs brought this "FTCA" action under 42 U.S.C. '
        'section 1983 - see at 1.', 13, cINK)
    notes(s, "OCR character errors are inherent to the scanned source — we remove "
             "STRUCTURAL noise (headers, stamps, page breaks) and repair unicode, but "
             "not OCR typos. Show the before/after. A separate aggressive normalisation "
             "(legal stop-words like 'plaintiff', 'court') feeds the TF-IDF features.")

    # ================= Slide 5 — Summarisation =================
    s = content_slide("Summarisation: long / short / tiny",
                      "Extractive — LexRank graph centrality over TF-IDF sentences")
    bullets(s, [
        "Clean → split into sentences → rank by graph centrality → take top-N per tier "
        "(long ≈ 12, short ≈ 4, tiny = 1).",
        "Filter out OCR/footnote fragments so the one-sentence ‘tiny’ stays meaningful.",
    ], top=Inches(1.7), height=Inches(1.3), size=16)
    # comparison cards (header text lives in the coloured bar so it stays visible)
    def cmp_card(l, w, hdr, hdr_col, body_runs):
        rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, Inches(3.2), w, Inches(2.95),
             fill=cWHITE, line=cRULE, line_w=Pt(1))
        bar = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, Inches(3.2), w, Inches(0.55),
                   fill=hdr_col)
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = bar.text_frame.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        run(bp, hdr, 15, cWHITE, bold=True)
        body = tbox(s, l + Inches(0.25), Inches(3.95), w - Inches(0.5), Inches(2.0))
        first = True
        for txt, sz, col, ital in body_runs:
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False; p.space_after = Pt(8)
            run(p, txt, sz, col, italic=ital)

    cmp_card(Inches(0.7), Inches(5.95), "OUR extractive tiny", cPRIM, [
        ("“…Katchen leasing agent Ms. Rodriguez signed the notice, acting as agent "
         "for Defendants Katchen and George and Helen Turk.”", 14, cINK, True)])
    cmp_card(Inches(6.9), Inches(5.7), "EXPERT tiny", cTEAL, [
        ("“Case regarding reasonable accommodation for service animals settles.”",
         14, cINK, True),
        ("Abstractive — paraphrases the whole case.", 12.5, cMUTED, False)])
    notes(s, "Extractive = we REUSE source sentences (no paraphrasing). LexRank builds a "
             "sentence-similarity graph and picks the most central sentences. Key fix: "
             "filter degenerate 'sentences' (footnote numbers, captions) so tiny is "
             "informative. Note the gap to the expert: ours captures a key fact, theirs "
             "abstracts the whole case — the extractive limitation.")

    # ================= Slide 6 — ROUGE =================
    s = content_slide("Do our summaries match the experts?",
                      "ROUGE F1 against the reference summaries")
    img(s, FIG / "rouge.png", Inches(1.35), Inches(1.7), w=Inches(10.6))
    bullets(s, [
        "ROUGE rises with tier length; LexRank > LSA > TextRank on long, noisy docs.",
        "Extractive can't paraphrase like experts → modest absolute ROUGE is expected.",
    ], top=Inches(5.85), height=Inches(1.2), size=14)
    notes(s, "ROUGE measures n-gram / LCS overlap with the expert summaries. Don't "
             "over-read absolute numbers: extractive summaries of 200-page docs won't "
             "match word-for-word, but the trend is the story — more budget = more "
             "overlap, and LexRank is the most robust summariser here.")

    # ================= Slide 7 — Modelling setup =================
    s = content_slide("Modelling setup", "Two tasks · four models · honest evaluation")
    bullets(s, [
        "Features: TF-IDF (1–2 grams, 30k vocab, sublinear TF) on the cleaned text.",
        "Targets: class_action_sought (binary) and case_type (12-class).",
        "Train on train+dev → report on the held-out TEST split.",
        "Metrics: overall accuracy, macro-F1, per-class accuracy, confusion matrix.",
    ], top=Inches(1.75), height=Inches(2.4), size=17)
    labels = [("Naive Bayes", "probabilistic", cPRIM),
              ("Logistic Reg.", "linear / GLM", cTEAL),
              ("Linear SVM", "max-margin", cCORAL),
              ("Bi-LSTM", "TensorFlow", cGOLD)]
    x = Inches(0.7)
    for name, sub, col in labels:
        card = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.5), Inches(2.9),
                    Inches(1.5), fill=cWHITE, line=cRULE, line_w=Pt(1))
        rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.5), Inches(2.9),
             Inches(0.13), fill=col)
        tf = card.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, name, 16, cINK, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        run(p2, sub, 12, cMUTED)
        x += Inches(3.05)
    notes(s, "We span the model families taught in class: probabilistic (NB), linear "
             "GLM (logistic regression), max-margin (SVM) and a neural sequence model "
             "(LSTM). Class-weighting handles the Yes/No imbalance. Everything reported "
             "on a held-out test split — no peeking.")

    # ================= Slide 8 — Results =================
    s = content_slide("Results — overall accuracy",
                      "Linear models on TF-IDF win on this subset")
    img(s, FIG / "leaderboard.png", Inches(0.55), Inches(1.7), w=Inches(8.7))
    stat_card(s, Inches(9.55), Inches(1.95), Inches(3.2), Inches(1.55),
              "0.95", "class action · Linear SVM", cCORAL)
    stat_card(s, Inches(9.55), Inches(3.75), Inches(3.2), Inches(1.55),
              "0.91", "case type · Linear SVM", cPRIM)
    bullets(s, [
        "LSTM trails — long docs truncated to 400 tokens + only ~1.5k training cases.",
    ], top=Inches(5.95), height=Inches(1.0), size=14)
    notes(s, "Headline: SVM 0.95 / 0.91. The interesting result is that simple TF-IDF + "
             "linear models BEAT the LSTM here — with long documents and a modest "
             "dataset, sparse bag-of-words generalises better than a sequence model "
             "trained from scratch. Logistic regression is the app default (best "
             "calibrated probabilities).")

    # ================= Slide 9 — per-class + confusion =================
    s = content_slide("Accuracy by class & confusion",
                      "Where the case-type model wins and loses")
    img(s, FIG / "perclass.png", Inches(0.5), Inches(1.75), w=Inches(6.1))
    img(s, FIG / "confusion.png", Inches(6.95), Inches(1.6), h=Inches(5.3))
    notes(s, "The assignment asks for accuracy BY class. Distinctive areas (Immigration, "
             "Equal Employment) are nearly perfect; the model confuses genuinely "
             "overlapping categories — Jail vs Prison Conditions, Policing vs Criminal "
             "Justice. The confusion matrix makes those off-diagonal mix-ups visible.")

    # ================= Slide 10 — interpretability =================
    s = content_slide("Interpretability — top words & word clouds",
                      "TF-IDF + linear models read like a glossary per class")
    img(s, FIG / "wc_class_action.png", Inches(2.37), Inches(1.62), w=Inches(8.6))
    img(s, FIG / "wc_case_type.png", Inches(1.27), Inches(4.55), w=Inches(10.8))
    notes(s, "Why we trust the models: the highest-weighted features are exactly the "
             "right legal vocabulary. class_action=Yes is driven by 'class', 'similarly "
             "situated', 'class certification' — textbook class-action language. Each "
             "case type surfaces its own glossary. Interpretable AND accurate.")

    # ================= Slide 11 — disagreement =================
    s = content_slide("Where do the models disagree?",
                      "Naive Bayes vs Linear SVM on case_type")
    bullets(s, [
        "Disagreements cluster on OVERLAPPING areas — Prison vs Jail Conditions, "
        "Policing vs Criminal Justice.",
        "Naive Bayes assumes word independence → over-counts generic legal terms and "
        "leans to the larger/overlapping class.",
        "Margin-based SVM uses the whole weighted bag-of-words → separates near-duplicate "
        "categories more reliably.",
        "Disagreement concentrates exactly where categories are semantically adjacent — "
        "a sensible, explainable failure mode.",
    ])
    notes(s, "We took the test cases where NB and SVM predict different case types and "
             "looked at which class pairs they swap — always adjacent categories. The "
             "WHY: NB's independence assumption over-weights common terms; SVM's margin "
             "uses full feature interaction, drawing cleaner boundaries. See notebook 03.")

    # ================= Slide 12 — LSTM =================
    s = content_slide("The TensorFlow model",
                      "Bi-LSTM architecture & accuracy over epochs")
    img(s, FIG / "lstm.png", Inches(1.35), Inches(1.7), w=Inches(10.6))
    bullets(s, [
        "Embedding → Bidirectional LSTM → global max-pool → dense; early stopping on val.",
        "Val accuracy plateaus early: 400-token truncation drops most of each case, and "
        "~1.5k examples is little for a from-scratch sequence model.",
    ], top=Inches(5.85), height=Inches(1.2), size=14)
    notes(s, "We DID build the required TensorFlow model — a bidirectional LSTM. The "
             "accuracy/epoch curves show it learning then plateauing/overfitting. It's "
             "the weakest model here — an honest result: truncation + small data. Future "
             "work: a long-document model (e.g. LED) that reads the whole case.")

    # ================= Slide 13 — demo =================
    s = content_slide("Live demo — the interactive tool",
                      "streamlit run app/streamlit_app.py")
    bullets(s, [
        "Paste ANY real or fake case → three summaries + both predictions with "
        "probability bars.",
        "Switch classifier (LogReg / NB / SVM / LSTM) and summariser live.",
        "Load a real test case to compare prediction against the true label.",
    ], top=Inches(1.7), height=Inches(1.8), size=16)
    box = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.9), Inches(3.85),
               Inches(9.5), Inches(2.7), fill=cWHITE, line=cPRIM, line_w=Pt(1.5))
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.05), Inches(3.78), Inches(9.2),
         Inches(0.16), fill=cCORAL)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, "Synthetic prison class-action sample  →", 15, cMUTED)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
    run(p, "class_action_sought:  Yes  (p ≈ 0.99)", 19, cTEAL, bold=True)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    run(p, "case_type:  Prison Conditions", 19, cPRIM, bold=True)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(8)
    run(p, "▶  insert a screen-recording of the app here", 12, cMUTED, italic=True)
    notes(s, "LIVE DEMO. Run `streamlit run app/streamlit_app.py`. 1) Sidebar → 'Use "
             "synthetic sample case' → Analyze → show the three summary tabs + Yes / "
             "Prison Conditions with probability bars. 2) Load a real test case and "
             "compare to its true label. 3) Flip the classifier dropdown to show model "
             "differences live. Replace this box with a screen recording in the video.")

    # ================= Slide 14 — wrap-up =================
    s = content_slide("Wrap-up & extra techniques", "What we applied and what's next")
    bullets(s, [
        "Course techniques: cleaning/normalisation, TF-IDF n-grams + sublinear TF, "
        "class-weighting, graph-based extractive summarisation, ROUGE, confusion-matrix "
        "error analysis, word clouds, model-disagreement analysis, a TensorFlow LSTM.",
        "Results: SVM 0.95 (class action) / 0.91 (case type) — interpretable & deployed "
        "in a live tool.",
        "Limitations: representative subset · extractive (not abstractive) · OCR noise.",
        "Future work: abstractive long-document summariser (LED/Pegasus); full 9k "
        "dataset; hierarchical models that read whole cases.",
    ])
    notes(s, "Recap techniques to tick the 'other course techniques' box. Be honest "
             "about limitations — subset, extractive, OCR. Future work points at "
             "abstractive long-doc models. Thank the audience and invite questions. "
             "Keep the whole talk under 10 minutes.")

    prs.save(DECK)
    print(f"  saved deck -> {DECK}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    generate_figures()
    build_deck()
    print("Done.")
